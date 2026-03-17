#!/usr/bin/env python3
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/home/mike/collateral/psychologist-project/slides/psychologist-case-study-v1.pptx")

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BG = RGBColor(0xF9, 0xF8, 0xF6)
AMBER = RGBColor(0xC9, 0x97, 0x3A)
TEAL = RGBColor(0x4A, 0x9B, 0x8E)
INK = RGBColor(0x24, 0x26, 0x2A)
MUTED = RGBColor(0x67, 0x6F, 0x7D)
LINE = RGBColor(0xD9, 0xD6, 0xCE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def add_text(slide, left, top, width, height, text, size=20, bold=False, color=INK,
             font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_small_label(slide, left, top, width, text, color=MUTED):
    return add_text(slide, left, top, width, Inches(0.25), text.upper(), size=10, bold=True, color=color)


def add_metric_pill(slide, left, top, width, value, label):
    add_box(slide, left, top, width, Inches(0.55), WHITE, LINE)
    add_text(slide, left + Inches(0.14), top + Inches(0.08), width - Inches(0.2), Inches(0.22), value, size=18, bold=True, color=NAVY)
    add_text(slide, left + Inches(0.14), top + Inches(0.29), width - Inches(0.2), Inches(0.18), label.upper(), size=8, bold=True, color=MUTED)


def add_bullet_list(slide, left, top, width, bullets, size=18, color=INK, bullet_color=AMBER):
    y = top
    for bullet in bullets:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, y + Inches(0.09), Inches(0.11), Inches(0.11))
        circ.fill.solid()
        circ.fill.fore_color.rgb = bullet_color
        circ.line.color.rgb = bullet_color
        add_text(slide, left + Inches(0.18), y, width - Inches(0.18), Inches(0.45), bullet, size=size, color=color)
        y += Inches(0.58)


def add_pipeline(slide, left, top, labels):
    box_w = Inches(1.85)
    gap = Inches(0.18)
    for i, (title, subtitle) in enumerate(labels):
        x = left + i * (box_w + gap)
        add_box(slide, x, top, box_w, Inches(1.1), WHITE, LINE)
        icon = add_box(slide, x + Inches(0.12), top + Inches(0.16), Inches(0.34), Inches(0.34), AMBER, AMBER)
        add_text(slide, x + Inches(0.56), top + Inches(0.14), box_w - Inches(0.66), Inches(0.22), title, size=16, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.12), top + Inches(0.56), box_w - Inches(0.24), Inches(0.34), subtitle, size=9, color=MUTED)
        if i < len(labels) - 1:
            add_text(slide, x + box_w + Inches(0.03), top + Inches(0.37), gap, Inches(0.2), "→", size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def add_chart(slide, left, top, width, height):
    add_box(slide, left, top, width, height, WHITE, LINE)
    chart_left = left + Inches(0.28)
    chart_top = top + Inches(0.38)
    chart_w = width - Inches(0.56)
    chart_h = height - Inches(0.9)
    # axes
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, chart_left, chart_top + chart_h, chart_w, Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.color.rgb = LINE
    line2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, chart_left, chart_top, Inches(0.02), chart_h)
    line2.fill.solid(); line2.fill.fore_color.rgb = LINE; line2.line.color.rgb = LINE
    vals = [0.22, 0.42, 0.56, 0.62, 0.79, 0.9]
    labels = ["T1", "T2", "T3", "T4", "T5", "T6"]
    points = []
    for i, v in enumerate(vals):
        x = chart_left + Inches(0.42) + i * ((chart_w - Inches(0.84)) / (len(vals) - 1))
        y = chart_top + chart_h - v * (chart_h - Inches(0.15))
        points.append((x, y))
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - Inches(0.05), y - Inches(0.05), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = AMBER; dot.line.color.rgb = AMBER
        add_text(slide, x - Inches(0.14), chart_top + chart_h + Inches(0.06), Inches(0.28), Inches(0.18), labels[i], size=9, color=MUTED, align=PP_ALIGN.CENTER)
    for (x1, y1), (x2, y2) in zip(points, points[1:]):
        connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
        connector.line.color.rgb = TEAL
        connector.line.width = Pt(2.5)
    add_small_label(slide, left + Inches(0.2), top + Inches(0.12), Inches(2.2), "Synthetic longitudinal trend")


def add_report_mock(slide, left, top, width, height, title, subtitle, lang):
    add_box(slide, left, top, width, height, WHITE, LINE)
    add_box(slide, left + Inches(0.16), top + Inches(0.14), width - Inches(0.32), Inches(0.32), BG, BG)
    add_text(slide, left + Inches(0.26), top + Inches(0.17), width - Inches(0.5), Inches(0.18), title, size=12, bold=True, color=NAVY)
    add_text(slide, left + Inches(0.26), top + Inches(0.4), width - Inches(0.5), Inches(0.16), subtitle, size=8, color=MUTED)
    add_box(slide, left + Inches(0.22), top + Inches(0.65), width - Inches(0.44), Inches(0.65), RGBColor(0xF4, 0xF7, 0xFA), LINE)
    add_text(slide, left + Inches(0.34), top + Inches(0.74), width - Inches(0.68), Inches(0.16), f"Patient summary ({lang})", size=10, bold=True, color=INK)
    add_text(slide, left + Inches(0.34), top + Inches(0.94), width - Inches(0.68), Inches(0.26), "Trend interpretation written for a non-technical reader.", size=9, color=MUTED)
    for i in range(3):
        add_box(slide, left + Inches(0.22), top + Inches(1.42) + i * Inches(0.34), width - Inches(0.44), Inches(0.12), RGBColor(0xE9, 0xEC, 0xF1), RGBColor(0xE9, 0xEC, 0xF1), radius=False)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_box(s, Inches(0.45), Inches(0.45), Inches(5.1), Inches(6.6), NAVY, NAVY)
    add_text(s, Inches(0.82), Inches(1.0), Inches(4.1), Inches(1.2), "Automated Clinical Reporting at Scale", size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.82), Inches(2.05), Inches(4.1), Inches(0.8), "From raw test data to multilingual patient reports — built for a practicing psychologist", size=15, color=RGBColor(0xE8, 0xEA, 0xEE))
    add_box(s, Inches(0.82), Inches(5.95), Inches(4.1), Inches(0.6), RGBColor(0x2A, 0x3C, 0x63), RGBColor(0x2A, 0x3C, 0x63))
    add_text(s, Inches(1.0), Inches(6.1), Inches(3.7), Inches(0.25), "EXPERT-VETTED · 100% JSS · 1,900+ HRS · MSEE GEORGIA TECH", size=9, bold=True, color=WHITE)
    add_report_mock(s, Inches(6.2), Inches(0.85), Inches(6.35), Inches(5.5), "Patient Progress Report", "Synthetic public-safe example", "Spanish")
    add_text(s, Inches(6.45), Inches(6.55), Inches(5.7), Inches(0.3), "Built with synthetic examples to protect client confidentiality.", size=10, color=MUTED)

    # Slide 2
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_text(s, Inches(0.7), Inches(0.55), Inches(6.2), Inches(0.55), "A Psychologist. Hundreds of Patients. Dozens of Test Types.", size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(1.12), Inches(4.2), Inches(0.3), "Manual reporting was the bottleneck.", size=16, color=MUTED)
    for i, (title, body) in enumerate([
        ("Growing patient volume", "No scalable reporting workflow for expanding caseload."),
        ("Hours lost per patient", "Manual formatting and interpretation consumed clinical time."),
        ("Multilingual delivery", "Patients needed results in language they could actually use."),
    ]):
        x = Inches(0.8 + i * 4.0)
        add_box(s, x, Inches(2.0), Inches(3.45), Inches(2.2), WHITE, LINE)
        add_box(s, x + Inches(0.18), Inches(2.2), Inches(0.38), Inches(0.38), AMBER, AMBER)
        add_text(s, x + Inches(0.72), Inches(2.15), Inches(2.35), Inches(0.28), title, size=16, bold=True, color=NAVY)
        add_text(s, x + Inches(0.18), Inches(2.72), Inches(3.0), Inches(0.8), body, size=11, color=MUTED)
    add_box(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.75), RGBColor(0xFA, 0xF1, 0xDE), RGBColor(0xFA, 0xF1, 0xDE))
    add_text(s, Inches(1.05), Inches(5.63), Inches(11.0), Inches(0.26), "The goal wasn't software. It was time back for the clinician.", size=18, bold=True, color=AMBER)

    # Slide 3
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_text(s, Inches(0.7), Inches(0.55), Inches(5.6), Inches(0.55), "A Full-Stack Reporting Engine", size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(1.12), Inches(5.6), Inches(0.3), "From data intake through multilingual PDF delivery", size=16, color=MUTED)
    add_pipeline(s, Inches(0.75), Inches(2.0), [
        ("Test Data Upload", "Bring repeated assessments into one workflow."),
        ("Longitudinal Analysis", "Show change over time, not snapshots."),
        ("Chart Generation", "Produce visuals that support interpretation."),
        ("Translation Layer", "Prepare summaries in the patient's language."),
        ("Final Report PDF", "Assemble a polished, clinician-ready document."),
    ])
    add_box(s, Inches(0.78), Inches(5.5), Inches(11.7), Inches(0.55), RGBColor(0xEE, 0xF3, 0xF6), RGBColor(0xEE, 0xF3, 0xF6))
    add_text(s, Inches(1.0), Inches(5.68), Inches(11.2), Inches(0.22), "Python · Flask · pandas · matplotlib · PostgreSQL · document generation", size=11, color=MUTED)

    # Slide 4
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_text(s, Inches(0.7), Inches(0.55), Inches(4.9), Inches(0.55), "What the Output Looks Like", size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(1.12), Inches(5.4), Inches(0.3), "Patient-friendly. Clinician-ready. Multilingual.", size=16, color=MUTED)
    add_chart(s, Inches(0.75), Inches(1.85), Inches(5.7), Inches(4.55))
    add_report_mock(s, Inches(6.75), Inches(1.85), Inches(5.8), Inches(4.55), "Patient Summary", "Synthetic bilingual example", "English / Spanish")
    add_box(s, Inches(8.2), Inches(5.95), Inches(3.4), Inches(0.62), RGBColor(0xFA, 0xF1, 0xDE), RGBColor(0xFA, 0xF1, 0xDE))
    add_text(s, Inches(8.45), Inches(6.13), Inches(3.0), Inches(0.22), "Generated in under 30 seconds", size=13, bold=True, color=AMBER)

    # Slide 5
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_text(s, Inches(0.7), Inches(0.55), Inches(4.9), Inches(0.55), "What Clients Say", size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(1.12), Inches(4.9), Inches(0.3), "Consistent quality. Long-term trust.", size=16, color=MUTED)
    quote = add_box(s, Inches(0.95), Inches(1.9), Inches(11.4), Inches(2.3), WHITE, LINE)
    add_text(s, Inches(1.3), Inches(2.3), Inches(10.7), Inches(0.95),
             "\"I can highly recommend Michael for your software development needs. Michael has been consistent in delivering quality work.\"",
             size=24, bold=True, color=NAVY)
    add_text(s, Inches(1.3), Inches(3.4), Inches(5.4), Inches(0.25), "VERIFIED UPWORK CLIENT · LONG-RUNNING SOFTWARE DEVELOPMENT ENGAGEMENT", size=10, bold=True, color=MUTED)
    add_metric_pill(s, Inches(1.1), Inches(4.95), Inches(3.2), "100% Job Success", "Upwork")
    add_metric_pill(s, Inches(4.6), Inches(4.95), Inches(3.2), "1,900+ Hours", "Billed")
    add_metric_pill(s, Inches(8.1), Inches(4.95), Inches(3.2), "Expert-Vetted", "By Upwork")

    # Slide 6
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_box(s, Inches(0.45), Inches(0.45), Inches(12.4), Inches(6.55), WHITE, LINE)
    add_text(s, Inches(0.95), Inches(0.9), Inches(5.0), Inches(0.5), "Is Your Problem Similar?", size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.95), Inches(1.45), Inches(5.6), Inches(0.5), "Complex data. Meaningful output. People who need clear answers.", size=16, color=MUTED)
    add_small_label(s, Inches(0.98), Inches(2.2), Inches(3.0), "Strong fit if you need")
    add_bullet_list(s, Inches(1.0), Inches(2.5), Inches(4.8), [
        "Large or growing data volume",
        "Reporting that needs to scale",
        "Output that non-technical users will actually read",
    ], size=18)
    add_small_label(s, Inches(7.0), Inches(2.2), Inches(3.0), "What this deck demonstrates")
    add_bullet_list(s, Inches(7.0), Inches(2.5), Inches(4.8), [
        "Python delivery with real workflow structure",
        "Analytics carried through into final reporting",
        "High-trust client work sustained over time",
    ], size=18, bullet_color=TEAL)
    add_box(s, Inches(0.98), Inches(5.65), Inches(11.0), Inches(0.72), RGBColor(0xF2, 0xF5, 0xF8), RGBColor(0xF2, 0xF5, 0xF8))
    add_text(s, Inches(1.2), Inches(5.9), Inches(10.6), Inches(0.24), "Upwork-safe version: no contact details included. This deck is intended as hiring collateral and project proof.", size=11, color=MUTED)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
